import PyPDF2
from docx import Document
from pptx import Presentation
from transformers import AutoProcessor, AutoModelForCausalLM
from PIL import Image
import torch
from google import genai
from google.genai import types
import cv2
import tempfile
from sentence_transformers import SentenceTransformer
import faiss

class Reader:
    def pdf_reader(self,uploaded_file):
        reader = PyPDF2.PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            txt = page.extract_text()
            if txt:
                text += txt + "\n"
        return text
     
    def docx_reader(self,uploaded_file):
         doc = Document(uploaded_file)
         text = "\n".join([para.text for para in doc.paragraphs])
         return text
     
    def text_reader(self,uploaded_file):
         text=uploaded_file.read().decode("utf-8", errors="ignore")
         return text
    
    def pptx_reader(self,uploaded_file):
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def image_reader(self,uploaded_file):
       
        image_bytes = uploaded_file.read()

        client = genai.Client(api_key="AIzaSyBxmzBpmBHR4B2Ym-mxc5FwM5mZRDaw6p0")
        response = client.models.generate_content(
            model='gemini-2.5-flash-lite',
            contents=[
            types.Part.from_bytes(
                data=image_bytes,
                mime_type='image/jpeg',
            ),
            'Extract all the text from the image'])
        return response.text
    
    def video_reader(self,uploaded_file):
        video_bytes = uploaded_file.read()

        client = genai.Client()
        response = client.models.generate_content(
            model='models/gemini-2.5-flash',
            contents=types.Content(
                parts=[
                    types.Part(
                        inline_data=types.Blob(data=video_bytes, mime_type='video/mp4')
                    ),
                    types.Part(text='Convert the video to text considering video frames and audio')
                ]
            )
        )
        return response.text

    def audio_reader(self,uploaded_file):
        audio_bytes = uploaded_file.read()
        client = genai.Client(api_key="AIzaSyBxmzBpmBHR4B2Ym-mxc5FwM5mZRDaw6p0")
        response = client.models.generate_content(
        model='gemini-2.5-flash',
        contents=[
            'extract text from the audio',
            types.Part.from_bytes(
            data=audio_bytes,
            mime_type='audio/mp3',
            )
        ]
        )
        return response.text
    
class QA:
    def __init__(self, text):
        self.text = text
        self.model = SentenceTransformer('all-MiniLM-L6-v2')
        self.chunks = self.splitter()
        self.index = self.create_faiss_index(self.chunks)
        self.client = genai.Client(api_key="AIzaSyBxmzBpmBHR4B2Ym-mxc5FwM5mZRDaw6p0")

    def splitter(self, chunk_size=500, overlap=200):
        """
        Split text into overlapping chunks for embedding and retrieval.
        """
        words = self.text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            chunks.append(chunk)
        return chunks

    def create_faiss_index(self, chunks):
        """
        Create a FAISS index for semantic search over text chunks.
        """
        embeddings = self.model.encode(chunks, normalize_embeddings=True)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)
        return index

    def QuestionAnswering(self, question):
        """
        Retrieve relevant context and generate an answer using Gemini.
        """
        query_embedding = self.model.encode([question], normalize_embeddings=True)
        D, I = self.index.search(query_embedding, k=3)
        context = "\n".join([self.chunks[i] for i in I[0]])
        prompt = f"""
        You are a helpful assistant. Use the provided context to answer the question accurately.

        Context:
        {context}

        Question:
        {question}

        Answer:
        """

        response = self.client.models.generate_content(
            model="gemini-2.5-flash-lite",
            contents=prompt
        )

        return response.text
